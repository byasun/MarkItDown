import os
import re
from markitdown import MarkItDown
from src.config import INPUT_DIR, OUTPUT_DIR
from src.utils import ensure_directory_exists, log_message, clean_text
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from pptx import Presentation

class Converter:
    def __init__(self):
        self.md = MarkItDown()
        ensure_directory_exists(INPUT_DIR)
        ensure_directory_exists(OUTPUT_DIR)

    def extract_images_from_pptx(self, pptx_path, img_output_dir):
        """Extrai imagens de um arquivo .pptx e as salva no diretório especificado."""
        ensure_directory_exists(img_output_dir)
        presentation = Presentation(pptx_path)
        img_paths = []

        for i, slide in enumerate(presentation.slides):
            for j, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # 13 = imagem
                    img_path = os.path.join(img_output_dir, f"slide_{i+1}_img_{j+1}.png")
                    with open(img_path, "wb") as f:
                        f.write(shape.image.blob)
                    img_paths.append(img_path)

        return img_paths

    def save_as_pdf(self, text, output_path):
        """Salva texto extraído em um PDF formatado (sem imagens)."""
        try:
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter  # Obtém tamanho da página
            margin = inch  # Margem para o texto
            max_width = width - 2 * margin  # Largura máxima para o texto

            # Configura fonte e posição inicial
            c.setFont("Helvetica", 12)
            y_position = height - inch  # Começa 1 polegada abaixo do topo

            # Quebra o texto em múltiplas linhas se for necessário
            lines = self.split_text(text, max_width, c)

            # Escreve cada linha no PDF
            for line in lines:
                if y_position < inch:  # Se atingir o final da página, cria uma nova
                    c.showPage()
                    c.setFont("Helvetica", 12)
                    y_position = height - inch
                c.drawString(margin, y_position, line)
                y_position -= 14  # Move para a próxima linha (14 pontos de espaçamento)

            c.save()
            log_message(f"PDF salvo: {output_path}")
        except Exception as e:
            log_message(f"Erro ao salvar PDF {output_path}: {e}", error=True)

    def split_text(self, text, max_width, canvas):
        """Quebra o texto em linhas de largura máxima definida (sem ultrapassar a largura da página)."""
        lines = []
        paragraphs = text.split('\n')  # Quebra o texto por parágrafos (quebra de linha)
        
        for paragraph in paragraphs:
            words = paragraph.split()  # Quebra o parágrafo em palavras
            current_line = ""

            for word in words:
                # Tenta adicionar a palavra à linha atual
                test_line = f"{current_line} {word}" if current_line else word
                text_width = canvas.stringWidth(test_line, "Helvetica", 12)

                # Se a largura da linha exceder a largura máxima, cria uma nova linha
                if text_width <= max_width:
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word

            # Adiciona a última linha
            if current_line:
                lines.append(current_line)

            # Adiciona uma linha em branco após cada parágrafo
            lines.append("")  

        return lines

    def process_files(self):
        files = os.listdir(INPUT_DIR)
        if not files:
            log_message("Nenhum arquivo encontrado na pasta de entrada.")
            return
        
        for file in files:
            input_path = os.path.join(INPUT_DIR, file)
            output_path = os.path.join(OUTPUT_DIR, f"{os.path.splitext(file)[0]}.pdf")
            img_output_dir = os.path.join(OUTPUT_DIR, f"{os.path.splitext(file)[0]}_images")

            try:
                result = self.md.convert(input_path)
                text_cleaned = clean_text(result.text_content)

                # Se for um arquivo .pptx, extrair imagens e salvar na pasta (sem adicioná-las no PDF)
                images = []
                if file.endswith(".pptx"):
                    images = self.extract_images_from_pptx(input_path, img_output_dir)

                # Exibir conteúdo extraído antes de salvar
                print(f"\n--- Conteúdo extraído de {file} ---\n")
                print(text_cleaned)
                print("\n------------------------------------\n")

                # Salva o texto no PDF sem imagens
                self.save_as_pdf(text_cleaned, output_path)
            except Exception as e:
                log_message(f"Erro ao converter {file}: {e}", error=True)
